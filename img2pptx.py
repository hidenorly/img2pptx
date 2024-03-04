#   Copyright 2023 hidenorly
#
#   Licensed under the Apache License, Version 2.0 (the "License");
#   you may not use this file except in compliance with the License.
#   You may obtain a copy of the License at
#
#       http://www.apache.org/licenses/LICENSE-2.0
#
#   Unless required by applicable law or agreed to in writing, software
#   distributed under the License is distributed on an "AS IS" BASIS,
#   WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#   See the License for the specific language governing permissions and
#   limitations under the License.

import argparse
import os
import collections.abc

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.dml.color import RGBColor

import webcolors
from PIL import Image
import pyheif

class ImageUtil:
    def getFilenameWithExt(filename, ext=".jpeg"):
        filename = os.path.splitext(filename)[0]
        return filename + ext

    def covertToJpeg(imageFile):
        outFilename = ImageUtil.getFilenameWithExt(imageFile, ".jpeg")
        image = None
        if imageFile.endswith(('.heic', '.HEIC')):
            heifImage = pyheif.read(imageFile)
            image = Image.frombytes(
                heifImage.mode,
                heifImage.size,
                heifImage.data,
                "raw",
                heifImage.mode,
                heifImage.stride,
            )
        else:
            image = Image.open(imageFile)
        if image:
            image.save(outFilename, "JPEG")
        return outFilename

class PowerPointUtil:
    SLIDE_WIDTH_INCH = 16
    SLIDE_HEIGHT_INCH = 9

    def __init__(self, path):
        self.prs = Presentation()
        self.prs.slide_width  = Inches(self.SLIDE_WIDTH_INCH)
        self.prs.slide_height = Inches(self.SLIDE_HEIGHT_INCH)
        self.path = path

    def save(self):
        self.prs.save(self.path)

    # layout is full, left, right, top, bottom
    def getLayoutPosition(self, layout="full"):
        # for full
        x=0
        y=0
        width = self.prs.slide_width
        height = self.prs.slide_height

        if layout=="left" or layout=="right":
            width = width /2
        if layout=="top" or layout=="bottom":
            height = height /2
        if layout=="right":
            x=width
        if layout=="bottom":
            y=height

        return x,y,width,height

    def getLayoutToFitRegion(self, width, height, regionWidth, regionHeight):
        resultWidth = width
        resultHeight = height

        if width > height:
            resultWidth = regionWidth
            resultHeight = int(regionWidth * height / width+0.99)
        else:
            resultHeight = regionHeight
            resultWidth = int(regionHeight * width / height+0.99)

        return resultWidth, regionHeight


    def addSlide(self, layout=None):
        if layout == None:
            layout = self.prs.slide_layouts[6]
        self.currentSlide = self.prs.slides.add_slide(layout)

    def addPicture(self, imagePath, x=0, y=0, width=None, height=None, isFitToSlide=True, regionWidth=None, regionHeight=None, isFitWihthinRegion=False, isCenter=False):
        if not regionWidth:
            regionWidth = self.prs.slide_width
        if not regionHeight:
            regionHeight = self.prs.slide_height
        regionWidth = int(regionWidth+0.99)
        regionHeight = int(regionHeight+0.99)
        pic = None
        try:
            pic = self.currentSlide.shapes.add_picture(imagePath, x, y)
        except:
            pass
        if pic:
            if width and height:
                pic.width = width
                pic.height = height
            else:
                if isFitToSlide:
                    width, height = pic.image.size
                    picWidth = pic.width
                    picHeight = pic.height
                    if width > height:
                        picWidth = regionWidth
                        picHeight = int(regionWidth * height / width + 0.99)
                    else:
                        picHeight = regionHeight
                        picWidth = int(regionHeight * width / height + 0.99)
                    if isFitWihthinRegion:
                        deltaWidth = picWidth - regionWidth
                        deltaHeight = picHeight - regionHeight
                        if deltaWidth>0 or deltaHeight>0:
                            # exceed the region
                            if deltaWidth > deltaHeight:
                                picWidth = regionWidth
                                picHeight = int(regionWidth * height / width + 0.99)
                            else:
                                picHeight = regionHeight
                                picWidth = int(regionHeight * width / height + 0.99)
                    if isCenter:
                        if picWidth != regionWidth:
                            pic.left = int( (regionWidth - picWidth) / 2 + 0.99 )
                        if picHeight != regionHeight:
                            pic.top = int( (regionHeight - picHeight) / 2 + 0.99 )

                    pic.width = picWidth
                    pic.height = picHeight
        return pic

    def nameToRgb(name):
        result = RGBColor(0,0,0)
        try:
            rgb = webcolors.name_to_rgb(name)
            result = RGBColor(rgb.red, rgb.green, rgb.blue)
        except:
            pass
        return result

    def applyExFormat(exFormat, textbox, font, text_frame):
        exFormats = exFormat.split(",")
        for anFormat in exFormats:
            cmdarg = anFormat.split(":")
            cmd = cmdarg[0]
            val = None
            if len(cmdarg)>=2:
                val = cmdarg[1]
            if cmd=="color":
                font.color.rgb = PowerPointUtil.nameToRgb(val)
            elif cmd=="face":
                font.name = val
            elif cmd=="size":
                font.size = Pt(float(val))
            elif cmd=="bold":
                font.bold = True
            elif cmd=="effect":
                # TODO: fix
                shadow = textbox.shadow
                shadow.visible = True
                shadow.shadow_type = 'outer'
                shadow.style = 'outer'
                shadow.blur_radius = Pt(5)
                shadow.distance = Pt(2)
                shadow.angle = 45
                shadow.color = MSO_THEME_COLOR_INDEX.ACCENT_5
                shadow.transparency = 0

    def addText(self, text, x=Inches(0), y=Inches(0), width=None, height=None, fontFace='Calibri', fontSize=Pt(18), isAdjustSize=True, textAlign = PP_ALIGN.LEFT, isVerticalCenter=False, exFormat=None):
        if width==None:
            width=self.prs.slide_width
        if height==None:
            height=self.prs.slide_height
        width = int(width+0.99)
        height = int(height+0.99)

        textbox = self.currentSlide.shapes.add_textbox(x, y, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        font = text_frame.paragraphs[0].font
        font.name = fontFace
        font.size = fontSize
        theHeight = textbox.height

        if exFormat:
            PowerPointUtil.applyExFormat(exFormat, textbox, font, text_frame)
        
        if isAdjustSize:
            text_frame.auto_size = True
            textbox.top = y

        if isVerticalCenter:
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        for paragraph in text_frame.paragraphs:
            paragraph.alignment = textAlign


if __name__=="__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("-i", "--input", default=".", help="Input folder path")
    parser.add_argument("-o", "--output", default="output.pptx", help="Output PowerPoint file path")
    parser.add_argument("-a", "--addFilename", default=False, action='store_true', help="Add filename to the slide")
    parser.add_argument("-l", "--layout", action='store', default='full', help='Specify layout full or left or right or center')
    parser.add_argument("-f", "--fullfit", action='store_true', default=False, help='Specify if want to fit within the slide')
    parser.add_argument('--offsetX', type=float, default=0, help='Specify offset x (Inch. max 16. float)')
    parser.add_argument('--offsetY', type=float, default=0, help='Specify offset y (Inch. max 9. float)')
    parser.add_argument('--fontFace', type=str, default="Calibri", help='Specify font face if necessary')
    parser.add_argument('--fontSize', type=float, default=18.0, help='Specify font size (pt) if necessary')
    parser.add_argument('--title', type=str, default=None, help='Specify title if necessary')
    parser.add_argument('--titleSize', type=float, default=None, help='Specify title size if necessary')
    parser.add_argument('--titleFormat', type=str, default=None, help='Specify title format if necessary e.g. color:white,face:Calibri,size:40,bold')
    args = parser.parse_args()

    prs = PowerPointUtil( args.output )

    imgPaths = []
    for dirpath, dirnames, filenames in os.walk(args.input):
        for filename in filenames:
            # convert required image
            if filename.endswith(('.heic', '.HEIC')):
                filename = ImageUtil.covertToJpeg(os.path.join(dirpath, filename))

            if filename.endswith(('.png', '.jpg', '.jpeg', '.JPG')):
                imgPaths.append( os.path.join(dirpath, filename) )

    imgPaths.sort()


    # --- add image file to the slide
    x, y, regionWidth, regionHeight = prs.getLayoutPosition(args.layout)
    offsetX = Inches(args.offsetX)
    offsetY = Inches(args.offsetY)
    x = x + offsetX
    y = y + offsetY
    regionWidth = int( regionWidth - offsetX )
    regionHeight = int( regionHeight - offsetY )

    isFitWihthinRegion = args.fullfit

    fontFace = args.fontFace
    fontSize = Pt(args.fontSize)

    textAlign = PP_ALIGN.LEFT
    if args.layout == "right":
        textAlign = PP_ALIGN.RIGHT

    titleSize = args.offsetY*72.0 #Inch to Pt
    if args.titleSize:
        titleSize = Pt(args.titleSize)
    if titleSize<100 or titleSize>400000:
        titleSize = Pt(40) # fail safe

    titleHeight = offsetY
    if titleHeight==0:
        titleHeight = titleSize

    isCenter = True if args.layout == "center" else False

    for filename in imgPaths:
        prs.addSlide()
        pic = prs.addPicture(filename, x, y, None, None, True, regionWidth, regionHeight, isFitWihthinRegion, isCenter)
        # Add Title
        if args.title:
            prs.addText(args.title, x, 0, regionWidth, titleHeight, fontFace, titleSize, True, textAlign, True, args.titleFormat)
        # Add filename(URL) at bottom
        if pic and args.addFilename:
            text = os.path.basename(filename)
            prs.addText(text, x, int(y+regionHeight-Inches(0.4)), regionWidth, Inches(0.4), fontFace, fontSize, True, textAlign)

    # --- save the ppt file
    prs.save()
